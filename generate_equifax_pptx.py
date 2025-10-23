#!/usr/bin/env python3
"""
Generate Equifax Australia Value Creation Strategy PowerPoint Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Equifax color scheme
EQUIFAX_RED = RGBColor(196, 30, 58)  # #C41E3A
DARK_GRAY = RGBColor(45, 45, 45)     # #2D2D2D
MEDIUM_GRAY = RGBColor(85, 85, 85)   # #555555
LIGHT_GRAY = RGBColor(102, 102, 102) # #666666
BG_GRAY = RGBColor(245, 245, 245)    # #F5F5F5

def create_presentation():
    """Create the main presentation object"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add red accent bar on left
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.1), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = EQUIFAX_RED
    left_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Equifax Australia Value Creation Strategy"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY
    title_para.font.name = 'Montserrat'

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Value-based positioning through bundling, data leverage, and AI-enabled insight"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = MEDIUM_GRAY
    subtitle_para.font.name = 'Montserrat'

    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(4), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Prepared for: Equifax Australia"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = LIGHT_GRAY
    footer_para.font.name = 'Montserrat'

    # Add date
    date_box = slide.shapes.add_textbox(Inches(8), Inches(6.8), Inches(1.5), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "October 23, 2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = LIGHT_GRAY
    date_para.font.name = 'Montserrat'
    date_para.alignment = PP_ALIGN.RIGHT

def add_content_slide(prs, title, bullets, page_num):
    """Add a standard content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Add red header bar
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.1)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = EQUIFAX_RED
    header_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY
    title_para.font.name = 'Montserrat'

    # Add underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Inches(0.02)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = EQUIFAX_RED
    underline.line.fill.background()

    # Add bullets
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8.5), Inches(4.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Montserrat'
        p.space_after = Pt(12)

    # Add footer
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7), Inches(4), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Equifax Australia Value Creation Strategy"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = LIGHT_GRAY
    footer_para.font.name = 'Montserrat'

    # Add page number
    page_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = str(page_num)
    page_para = page_frame.paragraphs[0]
    page_para.font.size = Pt(12)
    page_para.font.color.rgb = LIGHT_GRAY
    page_para.font.name = 'Montserrat'
    page_para.alignment = PP_ALIGN.RIGHT

def main():
    """Generate the complete presentation"""
    print("Creating Equifax presentation...")
    prs = create_presentation()

    # Slide 1: Title
    add_title_slide(prs)
    print("✓ Slide 1: Title slide")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Equifax Australia is market leader with 16M+ consumers and 2M+ businesses in database",
        "Strategic objective: Shift to value-based positioning to increase WTP and defend price power",
        "Core patterns: Bundling, Lock-In (integration), Data Leverage, AI/Digitalization",
        "WTP model: WTP = α(IQS × BCS) + β(ID) + γ(CES) + δ(TS)",
        "Scenario: Bundled Enterprise drives ~2.4x higher WTP vs Basic in exemplar model (0.89 vs 0.37)",
        "Competitive context: Experian + illion (AUD 820M) scaling data breadth",
        "Roadmap: Q1–Q4 rollout of bundles, API stickiness, compliance dashboards, AI explainability",
        "KPIs: Bundle penetration, Integration depth, Compliance usage, Trust/NPS, WTP index"
    ], 2)
    print("✓ Slide 2: Executive Summary")

    # Slide 3: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Mission: Help people and organizations live their financial best by delivering trusted data insights",
        "Equifax Inc. global presence; Australia HQ in North Sydney",
        "16M+ Australian consumer records and 2M+ business profiles",
        "Part of global workforce of 14,000+ employees across 24 countries",
        "Credit Information & Risk Management: Comprehensive consumer credit reports and scores",
        "Fraud Prevention & Identity Verification: Advanced solutions like IDMatrix",
        "Commercial Insights & Analytics: SwiftCheck, financial viability assessments",
        "Data-Driven Marketing: Targeted customer engagement using proprietary datasets"
    ], 3)
    print("✓ Slide 3: Company Overview")

    # Slide 4: Industry & Competitive Context
    add_content_slide(prs, "Industry & Competitive Context", [
        "Market: Credit reporting & data analytics (ANZSIC 7293) – high barriers, regulated",
        "Market concentration: 3 major bureaus; Equifax primary provider to most large lenders",
        "Comprehensive Credit Reporting (CCR) expands positive data use",
        "Tech shift: AI/ML, automation, real-time decisioning transforming credit scoring",
        "Heightened data privacy & cybersecurity expectations following global scrutiny",
        "Experian Australia: Global platform with consumer/business credit, fraud detection",
        "illion: Acquired by Experian for ~AUD 820M; formerly Dun & Bradstreet ANZ",
        "Indirect competitors: Credit Savvy, ClearScore, GetCreditScore, CreditorWatch"
    ], 4)
    print("✓ Slide 4: Industry & Competitive Context")

    # Slide 5: Strategic Challenge
    add_content_slide(prs, "Strategic Challenge – Value-Based Positioning", [
        "Objective: Increase customers' willingness-to-pay by elevating perceived value",
        "Goal: Expand bundle coverage across use cases (credit, ID, fraud, analytics)",
        "Goal: Deepen API integration to raise switching costs",
        "Goal: Simplify compliance and elevate trust",
        "Goal: Differentiate via AI insight quality and explainability",
        "Success Metrics: Premium price realization, Higher ARPU, Lower churn",
        "Success Metrics: Higher bundle penetration"
    ], 5)
    print("✓ Slide 5: Strategic Challenge")

    # Slide 6: Business Model Pattern Framework
    add_content_slide(prs, "Business Model Pattern Framework", [
        "Value Creation: Bundling, Lock-In, Data Leverage, AI/Digitalization → Increase WTP",
        "Pricing: Value-Based Pricing, Hybrid models, Price Discrimination → Capture WTP",
        "Power: Switching Costs, Network Effects, Brand/Reputation → Defend Margins",
        "Cost Efficiency: Economies of Scale, Vertical Integration, Platform-as-a-Service",
        "Four categories drive value across different dimensions",
        "Focus on WTP (Willingness-To-Pay) as primary value driver"
    ], 6)
    print("✓ Slide 6: Business Model Pattern Framework")

    # Slide 7: Value Creation Patterns
    add_content_slide(prs, "Value Creation Patterns – Detailed", [
        "Bundling: Credit + ID + Fraud + Analytics packages; tiered to roles and outcomes",
        "Lock-In (Integration): API embedding into underwriting/fraud workflows",
        "Data Leverage: Proprietary datasets; local market signals; client data co-development",
        "AI/Digitalization: Real-time ML scoring; explainability dashboards",
        "Process Flow: Data Ingestion → Analytics → Decision → Monitoring",
        "16M+ consumers, 2M+ businesses feeding ML models and risk scoring",
        "API delivery and workflow integration for seamless decision-making"
    ], 7)
    print("✓ Slide 7: Value Creation Patterns")

    # Slide 8: WTP Formula Introduction
    add_content_slide(prs, "WTP Formula Introduction", [
        "WTP = α(IQS × BCS) + β(ID) + γ(CES) + δ(TS)",
        "IQS: Insight Quality Score - Measures accuracy, timeliness, and depth",
        "BCS: Bundle Coverage Score - Breadth of use cases covered",
        "ID: Integration Depth - Degree of API/workflow embedding",
        "CES: Compliance Ease Score - Simplicity of regulatory compliance support",
        "TS: Trust Score - Brand perception, security confidence, reliability",
        "Weights α, β, γ, δ calibrated through customer research and pricing experiments"
    ], 8)
    print("✓ Slide 8: WTP Formula")

    # Slide 9: Operational Quantification
    add_content_slide(prs, "Operational Quantification – Variables & Scoring", [
        "IQS: Model AUC/KS, latency, coverage (0-1) | Target: 0.8-1.0",
        "BCS: Modules normalized (0-1) | Target: 0.7-1.0",
        "ID: Endpoints/workflows per client (0-1) | Target: 5+ endpoints",
        "CES: Customer ratings, audit feature usage (0-1) | Target: 0.75-1.0",
        "TS: NPS/trust index (0-1) | Target: 8/10+",
        "Data sources: Analytics metrics, client utilization, API usage, surveys",
        "All variables measured and tracked quarterly"
    ], 9)
    print("✓ Slide 9: Operational Quantification")

    # Slide 10: WTP Scenario Analysis
    add_content_slide(prs, "WTP Scenario Analysis – Tiers", [
        "Weights: α=0.5, β=0.2, γ=0.2, δ=0.1",
        "Basic Tier: IQS=0.60, BCS=0.50, ID=0.30, CES=0.40, TS=0.80 → WTP=0.37",
        "Pro Tier: IQS=0.75, BCS=0.70, ID=0.60, CES=0.70, TS=0.80 → WTP=0.60",
        "Enterprise Tier: IQS=0.90, BCS=1.00, ID=0.90, CES=0.90, TS=0.80 → WTP=0.89",
        "Enterprise tier shows ~2.4x higher WTP than Basic tier",
        "Bundle coverage and insight quality are key multiplicative drivers",
        "Integration depth and compliance ease add significant incremental value"
    ], 10)
    print("✓ Slide 10: WTP Scenario Analysis")

    # Slide 11: Competitive Benchmarking
    add_content_slide(prs, "Competitive Benchmarking Framework", [
        "Scope: Data breadth/depth, analytics sophistication, integration capabilities",
        "Equifax: Strong integration depth (0.8), solid insight quality (0.75)",
        "Experian+illion: Higher insight quality (0.85), good integration (0.65)",
        "Post-acquisition Experian+illion gaining ground in data breadth",
        "Equifax maintains lead in workflow integration and API embedding",
        "Strategic positioning: Double down on integration depth and compliance utility",
        "Focus on trust and regulatory compliance as differentiators"
    ], 11)
    print("✓ Slide 11: Competitive Benchmarking")

    # Slide 12: Benchmark Insights
    add_content_slide(prs, "Benchmark Insights – Variable-by-Variable", [
        "IQS: Leverage larger dataset and local market expertise for unique insights",
        "BCS: Differentiate by deepening niche modules like public-sector compliance",
        "ID: Use incumbent status to raise switching costs through deeper integrations",
        "CES: Emphasize regulatory domain expertise in Australia (CCR, Privacy Act)",
        "TS: Reinforce trust advantage through security certifications and transparency",
        "Opportunity: Build automated compliance dashboards with audit-ready reporting",
        "Opportunity: Promote 'unique insight' and 'local market expertise'"
    ], 12)
    print("✓ Slide 12: Benchmark Insights")

    # Slide 13: Strategic Implications
    add_content_slide(prs, "Strategic Implications", [
        "Push IQS × BCS: Develop broader, deeper bundles with superior insight quality",
        "Maximize Integration Depth: Implement co-development initiatives and plug-ins",
        "Elevate Compliance Ease: Deploy compliance dashboards and automated alerts",
        "Reinforce Trust Score: Showcase security certifications and regulatory leadership",
        "Value-Based Pricing: Implement tiered, outcome-linked pricing structures",
        "Focus on multiplicative value drivers (IQS × BCS) for maximum impact",
        "Build switching costs through deep API and workflow integration"
    ], 13)
    print("✓ Slide 13: Strategic Implications")

    # Slide 14: Implementation Roadmap
    add_content_slide(prs, "Implementation Roadmap Overview (Q1–Q4)", [
        "Q1: Finalize bundle structure and pricing tiers",
        "Q1-Q2: Launch initial bundles with pilot clients",
        "Q2-Q3: Develop enhanced API integration toolkit and developer portal",
        "Q2-Q4: Pilot and deploy compliance dashboard across client base",
        "Q3-Q4: Develop and deploy AI explainability modules",
        "Q1-Q4: Continuous client interviews and co-development programs",
        "Q4: Full rollout of enterprise bundles and integration capabilities"
    ], 14)
    print("✓ Slide 14: Implementation Roadmap")

    # Slide 15: Operational Levers - Bundling
    add_content_slide(prs, "Operational Levers – Bundling & Integration", [
        "Bundling: Tiered Solutions - Essentials, Advanced, and Enterprise packages",
        "Role-Based Packaging: Tailored to risk managers, compliance officers, fraud analysts",
        "Outcome Messaging: Frame bundles around measurable outcomes (reduce default risk)",
        "Integration: Custom API Toolkits with sandbox environments",
        "Decision Engine Connectors: Embed into underwriting and fraud workflows",
        "Co-Development Hooks: ML models customizable with client outcome data",
        "Control: Monitor cross-sell rates and API client engagement"
    ], 15)
    print("✓ Slide 15: Operational Levers - Bundling")

    # Slide 16: Operational Levers - Compliance
    add_content_slide(prs, "Operational Levers – Compliance & Trust", [
        "Compliance: RegTech-as-a-Service for CCR, AML/KYC, Privacy Act obligations",
        "Compliance Dashboards: Visualize audit trails and data usage",
        "Automated Alerts: Notifications when portfolios breach regulatory thresholds",
        "Trust: Cybersecurity Differentiation - Promote certifications and encryption",
        "Local Regulatory Alignment: Highlight role in shaping AU-specific regulations",
        "Case Studies: Develop proof points of fraud/loss prevention",
        "Transparency: Offer explainability tools for AI models in credit scoring"
    ], 16)
    print("✓ Slide 16: Operational Levers - Compliance")

    # Slide 17: Control Points & KPIs
    add_content_slide(prs, "Control Points & KPI Dashboard", [
        "Bundle Penetration Rate: Target 75%, Current 65%",
        "Average Integration Depth: Target 5 endpoints, Current 4.2",
        "Compliance Module Usage: Target 80%, Current 68%",
        "Trust/NPS Score: Target 8/10, Current 7.2",
        "WTP Index: Target 0.8, Current 0.65",
        "Measurement Frequency: Monthly for operational, Quarterly for strategic",
        "Control Points: Monthly leadership dashboard review, Quarterly bundle adjustments"
    ], 17)
    print("✓ Slide 17: Control Points & KPIs")

    # Slide 18: Risk Management
    add_content_slide(prs, "Risk Management Framework", [
        "Threat: Low-price competitor bundles (High Risk)",
        "Mitigation: Flexible modular pricing & highlight AI/compliance differentiation",
        "Threat: Economic downturn usage reduction (Medium-High Risk)",
        "Mitigation: Performance pricing and ROI guarantees",
        "Threat: Regulatory data access restrictions (Medium Risk)",
        "Mitigation: Analytics innovation using permitted data, diversify signal sources",
        "Threat: SME integration burden (Low-Medium Risk)",
        "Mitigation: Build no-code UI/API bridge for less technical clients"
    ], 18)
    print("✓ Slide 18: Risk Management")

    # Slide 19: WTP Variable Tracking
    add_content_slide(prs, "WTP Variable Tracking Template", [
        "IQS: Owner - Head of Data Science | Current 0.78 | Launch explainability module Q4",
        "BCS: Owner - Product Strategy Lead | Current 0.65 | Introduce compliance scoring",
        "ID: Owner - CTO | Current 4.2 endpoints | Expand developer portal & SDK",
        "CES: Owner - Risk & Legal | Current 0.68 | Add automated regulatory reports",
        "TS: Owner - Customer Success | Current 7.2 | Publish security certifications",
        "All variables measured quarterly with assigned owners",
        "Action items tracked to close gaps to target ranges"
    ], 19)
    print("✓ Slide 19: WTP Variable Tracking")

    # Slide 20: Key Recommendations
    add_content_slide(prs, "Key Recommendations", [
        "1. Accelerate bundling: Implement tiered packages with outcome-led offerings",
        "2. Deepen integration: Enhance API toolkits and embed into client workflows",
        "3. Launch compliance dashboards: Develop RegTech-as-a-Service solution",
        "4. Embed AI explainability: Create transparency tools for AI-driven decisions",
        "5. Establish WTP governance: Implement quarterly measurement of all variables",
        "Focus on raising switching costs through deeper integration",
        "Position Equifax as the compliance simplifier for Australian market"
    ], 20)
    print("✓ Slide 20: Key Recommendations")

    # Slide 21: Next Steps & Closing
    add_content_slide(prs, "Next Steps & Closing", [
        "Next 30-60 Days: Validate WTP Weights via 10-15 Client Interviews",
        "Finalize Bundle SKUs and Pricing - Launch pilot with 3 anchor clients",
        "Stand Up KPI Dashboard - Implement tracking and assign metric owners",
        "Prepare Security/Compliance Collateral - Develop trust-building materials",
        "Establish quarterly review cadence for WTP variable tracking",
        "Begin co-development programs with key enterprise clients",
        "Contact: Your Equifax Australia Strategy Team",
        "Thank You"
    ], 21)
    print("✓ Slide 21: Next Steps & Closing")

    # Save presentation
    filename = "Equifax_Australia_Value_Creation_Strategy.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation saved as: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    main()
